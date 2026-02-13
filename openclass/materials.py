"""
OpenClass 课堂材料解析模块
支持解析 PPT、PDF、Word 等文件，提取文本内容作为 AI 上下文
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


class MaterialParser:
    """课堂材料解析器"""

    @staticmethod
    def parse(file_path: str) -> str:
        """
        解析文件并提取文本内容
        支持: .pptx, .pdf, .docx, .txt, .md
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        suffix = path.suffix.lower()
        parsers = {
            ".pptx": MaterialParser._parse_pptx,
            ".pdf": MaterialParser._parse_pdf,
            ".docx": MaterialParser._parse_docx,
            ".doc": MaterialParser._parse_docx,
            ".txt": MaterialParser._parse_text,
            ".md": MaterialParser._parse_text,
            ".csv": MaterialParser._parse_text,
        }

        parser = parsers.get(suffix)
        if not parser:
            logger.warning(f"不支持的文件格式: {suffix}，尝试作为文本读取")
            return MaterialParser._parse_text(file_path)

        try:
            text = parser(file_path)
            logger.info(f"已解析材料: {path.name} ({len(text)} 字符)")
            return text
        except Exception as e:
            logger.error(f"解析材料失败: {file_path} - {e}")
            return ""

    @staticmethod
    def _parse_pptx(file_path: str) -> str:
        """解析 PowerPoint 文件"""
        from pptx import Presentation

        prs = Presentation(file_path)
        texts = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts = [f"--- 幻灯片 {slide_idx} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                # 处理表格
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            slide_texts.append(row_text)
            texts.append("\n".join(slide_texts))

        return "\n\n".join(texts)

    @staticmethod
    def _parse_pdf(file_path: str) -> str:
        """解析 PDF 文件"""
        from PyPDF2 import PdfReader

        reader = PdfReader(file_path)
        texts = []
        for page_idx, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                texts.append(f"--- 页面 {page_idx} ---\n{text.strip()}")

        return "\n\n".join(texts)

    @staticmethod
    def _parse_docx(file_path: str) -> str:
        """解析 Word 文件"""
        from docx import Document

        doc = Document(file_path)
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text.strip())

        # 处理表格
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    texts.append(row_text)

        return "\n".join(texts)

    @staticmethod
    def _parse_text(file_path: str) -> str:
        """解析纯文本文件"""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    @staticmethod
    def parse_multiple(file_paths: list[str]) -> str:
        """解析多个文件，合并文本"""
        all_texts = []
        for fp in file_paths:
            text = MaterialParser.parse(fp)
            if text:
                name = Path(fp).name
                all_texts.append(f"=== 材料: {name} ===\n{text}")
        return "\n\n".join(all_texts)
